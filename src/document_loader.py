"""
文档加载与分块模块

支持的文件格式：
  - .txt: 纯文本文件
  - .docx: Word文档
  - .xlsx: Excel表格
  - .pptx: PowerPoint演示文稿
  - .pdf: PDF文档
  - .html: HTML网页

文档级别支持：
  - 文档可以设置安全级别：public, internal, confidential, secret
  - 级别可以通过文件名前缀指定，如：[confidential]敏感文档.txt
  - 也可以通过配置文件 data/doc_levels.json 指定

为什么需要分块（Chunking）？
  - LLM 有上下文窗口限制，不能把整本书塞进一个 prompt
  - 检索时，细粒度的块比整篇文档更容易匹配到用户问题
  - 块与块之间保留重叠（overlap），避免关键信息恰好被切在边界上

工程要点：
  - 用 pathlib.Path 处理路径（跨平台，比 os.path 更现代）
  - 用 try/except 包裹 IO 操作，给用户清晰的错误信息
  - 按需导入库，避免不必要的依赖加载
"""

import json
from pathlib import Path
from typing import List, Dict, Optional

# 文档级别定义
DOC_LEVELS = ["public", "internal", "confidential", "secret"]
DEFAULT_DOC_LEVEL = "internal"  # 默认文档级别


def parse_doc_level(filename: str) -> tuple[str, str]:
    """
    从文件名中解析文档级别
    
    支持的格式：
      - [secret]敏感文档.txt -> ("secret", "敏感文档.txt")
      - [confidential]机密.txt -> ("confidential", "机密.txt")
      - 普通文档.txt -> ("internal", "普通文档.txt")
    
    Args:
        filename: 文件名
        
    Returns:
        (文档级别, 清理后的文件名)
    """
    # 检查是否有级别前缀
    if filename.startswith("[") and "]" in filename:
        end = filename.index("]")
        level = filename[1:end].strip().lower()
        if level in DOC_LEVELS:
            clean_name = filename[end+1:].strip()
            return level, clean_name
    
    return DEFAULT_DOC_LEVEL, filename


def load_doc_levels_config(config_path: str = "data/doc_levels.json") -> Dict[str, str]:
    """
    加载文档级别配置文件
    
    Args:
        config_path: 配置文件路径
        
    Returns:
        文件名 -> 级别的映射
    """
    path = Path(config_path)
    if path.exists():
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, FileNotFoundError):
            pass
    return {}


def get_doc_level(filename: str, config: Optional[Dict[str, str]] = None) -> str:
    """
    获取文档的级别
    
    优先级：
      1. 配置文件中的设置
      2. 文件名前缀
      3. 默认级别
    
    Args:
        filename: 文件名
        config: 文档级别配置
        
    Returns:
        文档级别
    """
    # 1. 检查配置文件
    if config and filename in config:
        level = config[filename].lower()
        if level in DOC_LEVELS:
            return level
    
    # 2. 检查文件名前缀
    level, _ = parse_doc_level(filename)
    return level


def load_text(file_path: str | Path) -> str:
    """读取 txt 文件内容"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文档不存在: {path}")
    return path.read_text(encoding="utf-8")


def load_docx(file_path: str | Path) -> str:
    """读取 Word 文档 (.docx)"""
    from docx import Document
    
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文档不存在: {path}")
    
    doc = Document(path)
    content = []
    for para in doc.paragraphs:
        content.append(para.text)
    return "\n".join(content)


def load_xlsx(file_path: str | Path) -> str:
    """读取 Excel 表格 (.xlsx)"""
    from openpyxl import load_workbook
    
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文档不存在: {path}")
    
    wb = load_workbook(path, read_only=True)
    content = []
    
    for sheet in wb.sheets:
        content.append(f"=== 工作表: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell else "" for cell in row)
            if row_text.strip():
                content.append(row_text)
    
    return "\n".join(content)


def load_pptx(file_path: str | Path) -> str:
    """读取 PowerPoint 演示文稿 (.pptx)"""
    from pptx import Presentation
    
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文档不存在: {path}")
    
    prs = Presentation(path)
    content = []
    
    for slide in prs.slides:
        content.append("=== 幻灯片 ===")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    
    return "\n".join(content)


def load_pdf(file_path: str | Path) -> str:
    """读取 PDF 文档"""
    from PyPDF2 import PdfReader
    
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文档不存在: {path}")
    
    reader = PdfReader(path)
    content = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            content.append(text)
    return "\n".join(content)


def load_html(file_path: str | Path) -> str:
    """读取 HTML 文件，提取纯文本内容"""
    from bs4 import BeautifulSoup
    
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文档不存在: {path}")
    
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    
    # 移除脚本和样式
    for script in soup(["script", "style"]):
        script.decompose()
    
    return soup.get_text(separator="\n", strip=True)


def load_file(file_path: str | Path) -> str:
    """
    根据文件扩展名自动选择合适的加载方法
    
    参数：
      file_path: 文件路径
    
    返回：
      文件的文本内容
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    
    loaders = {
        ".txt": load_text,
        ".docx": load_docx,
        ".xlsx": load_xlsx,
        ".pptx": load_pptx,
        ".pdf": load_pdf,
        ".html": load_html,
    }
    
    if ext not in loaders:
        raise ValueError(f"不支持的文件格式: {ext}")
    
    return loaders[ext](path)


def chunk_text(text: str, chunk_size: int, overlap: int) -> List[str]:
    """
    将文本分割成重叠的块

    参数：
      chunk_size: 每块的最大字符数
      overlap:    相邻块之间的重叠字符数

    返回：
      字符串列表，每个元素是一个文本块
    """
    if chunk_size <= overlap:
        raise ValueError("chunk_size 必须大于 overlap")

    chunks: List[str] = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = start + chunk_size

        if end >= text_len:
            chunks.append(text[start:])
            break

        search_start = max(start + chunk_size * 3 // 4, start)
        cut = -1
        for sep in ("。", "！", "？", "\n", ".", "!", "?"):
            pos = text.rfind(sep, search_start, end)
            if pos > cut:
                cut = pos

        if cut > search_start:
            chunks.append(text[start: cut + 1])
            start = cut + 1 - overlap
        else:
            chunks.append(text[start:end])
            start = end - overlap

    if start < 0:
        chunks[-1] = text

    return chunks


def load_and_chunk(file_path: str | Path, chunk_size: int = 500, overlap: int = 100) -> List[str]:
    """加载文档并分块的一站式函数"""
    text = load_file(file_path)
    chunks = chunk_text(text, chunk_size, overlap)
    return chunks


def load_documents_from_folder(folder_path: str | Path, chunk_size: int = 500, overlap: int = 100, specific_files: List[str] = None) -> tuple[List[str], dict, List[str]]:
    """
    加载文件夹中的所有支持格式的文档并分块

    支持的格式：.txt, .docx, .xlsx, .pptx, .pdf, .html

    参数：
      folder_path: 文件夹路径
      chunk_size: 每块的最大字符数
      overlap:    相邻块之间的重叠字符数
      specific_files: 可选，指定要加载的文件列表（用于增量更新）

    返回：
      (all_chunks, doc_metadata, chunk_sources) - 文档块列表、文档元数据、每个块的来源文件名（包含级别）
      
    注意：
      chunk_sources 的格式为 "filename|level"，例如 "report.txt|internal"
    """
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        raise FileNotFoundError(f"文档文件夹不存在: {folder}")

    all_chunks: List[str] = []
    chunk_sources: List[str] = []
    doc_metadata: dict = {}
    supported_extensions = ("*.txt", "*.docx", "*.xlsx", "*.pptx", "*.pdf", "*.html")
    
    # 加载文档级别配置
    doc_levels_config = load_doc_levels_config()
    
    all_files = []
    for ext in supported_extensions:
        all_files.extend(folder.glob(ext))
    
    # 如果指定了要加载的文件列表，过滤掉其他文件
    if specific_files:
        specific_set = set(specific_files)
        all_files = [f for f in all_files if f.name in specific_set]
    
    if not all_files:
        raise FileNotFoundError(f"文件夹中没有找到支持的文档格式: {folder}")

    for doc_file in all_files:
        print(f"  加载文档: {doc_file.name}")
        try:
            chunks = load_and_chunk(doc_file, chunk_size, overlap)
            all_chunks.extend(chunks)
            
            # 获取文档级别
            doc_level = get_doc_level(doc_file.name, doc_levels_config)
            
            # chunk_sources 格式: "filename|level"
            source_with_level = f"{doc_file.name}|{doc_level}"
            chunk_sources.extend([source_with_level] * len(chunks))
            
            # 记录文档元数据（包含级别）
            doc_metadata[doc_file.name] = {
                "mtime": doc_file.stat().st_mtime,
                "size": doc_file.stat().st_size,
                "level": doc_level
            }
            print(f"    -> 分割为 {len(chunks)} 个文本块 (级别: {doc_level})")
        except Exception as e:
            print(f"    -> 加载失败: {e}")

    return all_chunks, doc_metadata, chunk_sources